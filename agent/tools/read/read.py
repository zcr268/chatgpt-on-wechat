"""
Read tool - Read file contents
Supports text files, images (jpg, png, gif, webp), and PDF files
"""

import os
from typing import Dict, Any
from pathlib import Path

from agent.tools.base_tool import BaseTool, ToolResult
from agent.tools.utils.credentials import DENIED_MESSAGE, is_credential_path
from agent.tools.utils.file_state import note_read
from agent.tools.utils.truncate import truncate_head, format_size, DEFAULT_MAX_LINES, DEFAULT_MAX_BYTES
from common.utils import expand_path

EMPTY_FILE_NOTICE = "[File exists but is empty]"


def split_lines(content: str) -> list:
    """Split into lines, dropping the empty element after a trailing newline.

    Without this a file ending in "\\n" reports one line too many and renders a
    phantom numbered blank line at the end.
    """
    lines = content.split('\n')
    if lines and lines[-1] == '':
        lines.pop()
    return lines


PDF_MAX_PAGES_PER_READ = 20


def _parse_page_range(pages, total_pages: int):
    """Parse a PDF ``pages`` argument into an inclusive (first, last) pair.

    Accepts "5", "2-8" and the open-ended "10-". Defaults to the first
    PDF_MAX_PAGES_PER_READ pages so a large document is never parsed whole.
    """
    if total_pages <= 0:
        return 1, 0

    spec = str(pages or "").strip()
    if not spec:
        return 1, min(PDF_MAX_PAGES_PER_READ, total_pages)

    if '-' in spec:
        head, _, tail = spec.partition('-')
        head, tail = head.strip(), tail.strip()
    else:
        head = tail = spec.strip()

    try:
        first = int(head) if head else 1
        last = int(tail) if tail else total_pages
    except ValueError:
        raise ValueError(f'Invalid pages value "{spec}". Use e.g. "3", "1-5" or "10-".')

    if first < 1 or last < first:
        raise ValueError(f'Invalid pages range "{spec}" for a {total_pages}-page document.')
    if first > total_pages:
        raise ValueError(f'Page {first} is beyond the end of the document ({total_pages} pages).')

    last = min(last, total_pages, first + PDF_MAX_PAGES_PER_READ - 1)
    return first, last


def number_lines(text: str, start_line: int) -> str:
    """Prefix each line with its 1-indexed file line number as ``12|code``.

    The compact gutter (no padding) keeps the token cost low while letting the
    model line up read output with grep hits and ask for a precise offset.
    Callers must tell the model these prefixes are display-only.
    """
    return '\n'.join(
        f"{start_line + i}|{line}" for i, line in enumerate(text.split('\n'))
    )


class Read(BaseTool):
    """Tool for reading file contents"""
    
    name: str = "read"
    description: str = f"Read or inspect file contents. For text/PDF/Word/Excel/PPT files, returns content (truncated to {DEFAULT_MAX_LINES} lines or {DEFAULT_MAX_BYTES // 1024}KB). Each line is prefixed with its line number as `12|content` - these prefixes are display aids for locating lines and are NOT part of the file, so strip them whenever you reuse the content (in edit's oldText/newText, when writing it elsewhere, or when quoting it to the user). For images/videos/audio, returns metadata only (file info, size, type). Use offset/limit for large text files."
    
    params: dict = {
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Path to the file to read. IMPORTANT: Relative paths are based on workspace directory. To access files outside workspace, use absolute paths starting with ~ or /."
            },
            "offset": {
                "type": "integer",
                "description": "Line number to start reading from (1-indexed, optional). Use negative values to read from end (e.g. -20 for last 20 lines)"
            },
            "limit": {
                "type": "integer",
                "description": "Maximum number of lines to read (optional)"
            },
            "pages": {
                "type": "string",
                "description": f"PDF only: page range to extract, e.g. \"3\", \"1-5\" or \"10-\". Defaults to the first {PDF_MAX_PAGES_PER_READ} pages; at most {PDF_MAX_PAGES_PER_READ} pages are read per call."
            }
        },
        "required": ["path"]
    }
    
    def __init__(self, config: dict = None):
        self.config = config or {}
        self.cwd = self.config.get("cwd", os.getcwd())
        
        # File type categories
        self.image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.svg', '.ico'}
        self.video_extensions = {'.mp4', '.avi', '.mov', '.mkv', '.flv', '.wmv', '.webm', '.m4v'}
        self.audio_extensions = {'.mp3', '.wav', '.ogg', '.m4a', '.flac', '.aac', '.wma'}
        self.binary_extensions = {'.exe', '.dll', '.so', '.dylib', '.bin', '.dat', '.db', '.sqlite'}
        self.archive_extensions = {'.zip', '.tar', '.gz', '.rar', '.7z', '.bz2', '.xz'}
        self.pdf_extensions = {'.pdf'}
        self.office_extensions = {'.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx'}

        # Readable text formats (will be read with truncation)
        self.text_extensions = {
            '.txt', '.md', '.markdown', '.rst', '.log', '.csv', '.tsv', '.json', '.xml', '.yaml', '.yml',
            '.py', '.js', '.ts', '.java', '.c', '.cpp', '.h', '.hpp', '.go', '.rs', '.rb', '.php',
            '.html', '.css', '.scss', '.sass', '.less', '.vue', '.jsx', '.tsx',
            '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd',
            '.sql', '.r', '.m', '.swift', '.kt', '.scala', '.clj', '.erl', '.ex',
            '.dockerfile', '.makefile', '.cmake', '.gradle', '.properties', '.ini', '.conf', '.cfg',
        }
    
    def execute(self, args: Dict[str, Any]) -> ToolResult:
        """
        Execute file read operation
        
        :param args: Contains file path and optional offset/limit parameters
        :return: File content or error message
        """
        # Support 'location' as alias for 'path' (LLM may use it from skill listing)
        path = args.get("path", "") or args.get("location", "")
        path = path.strip() if isinstance(path, str) else ""
        offset = args.get("offset")
        limit = args.get("limit")
        pages = args.get("pages")

        if not path:
            return ToolResult.fail("Error: path parameter is required")
        
        # Resolve path
        absolute_path = self._resolve_path(path)
        
        # Security check: block credential files and their aliases.
        # See issue #2913 (/proc/self/environ bypass) and #2863 (scope).
        if self._is_credential_path(absolute_path):
            return ToolResult.fail(DENIED_MESSAGE)
        
        # Check if file exists. Both misses below name the tool that answers
        # the question instead - otherwise the model tends to guess at nearby
        # paths and burn turns on more failed reads.
        if not os.path.exists(absolute_path):
            hint = "Hint: use search_files with target='files' to locate it by name."
            # Provide helpful hint if using relative path
            if not os.path.isabs(path) and not path.startswith('~'):
                return ToolResult.fail(
                    f"Error: File not found: {path}\n"
                    f"Resolved to: {absolute_path}\n"
                    f"Hint: Relative paths are based on workspace ({self.cwd}). For files outside workspace, use absolute paths.\n"
                    f"{hint}"
                )
            return ToolResult.fail(f"Error: File not found: {path}\n{hint}")

        if os.path.isdir(absolute_path):
            return ToolResult.fail(
                f"Error: {path} is a directory, not a file. "
                f"Use the ls tool to list what is inside it."
            )
        
        # Check if readable
        if not os.access(absolute_path, os.R_OK):
            return ToolResult.fail(f"Error: File is not readable: {path}")
        
        # Check file type
        file_ext = Path(absolute_path).suffix.lower()
        file_size = os.path.getsize(absolute_path)
        
        # Check if image - return metadata for sending
        if file_ext in self.image_extensions:
            return self._read_image(absolute_path, file_ext)
        
        # Check if video/audio/binary/archive - return metadata only
        if file_ext in self.video_extensions:
            return self._return_file_metadata(absolute_path, "video", file_size)
        if file_ext in self.audio_extensions:
            return self._return_file_metadata(absolute_path, "audio", file_size)
        if file_ext in self.binary_extensions or file_ext in self.archive_extensions:
            return self._return_file_metadata(absolute_path, "binary", file_size)
        
        # Check if PDF
        if file_ext in self.pdf_extensions:
            return self._read_pdf(absolute_path, path, offset, limit, pages)

        # Check if Office document (.docx, .xlsx, .pptx, etc.)
        if file_ext in self.office_extensions:
            return self._read_office(absolute_path, path, file_ext, offset, limit)

        # Read text file (with truncation for large files)
        return self._read_text(absolute_path, path, offset, limit)
    
    def _resolve_path(self, path: str) -> str:
        """
        Resolve path to absolute path
        
        :param path: Relative or absolute path
        :return: Absolute path
        """
        # Expand ~ to user home directory
        path = expand_path(path)
        if os.path.isabs(path):
            return path
        return os.path.abspath(os.path.join(self.cwd, path))

    @staticmethod
    def _paginate(all_lines, offset, limit, display_path):
        """Slice, truncate and number a list of lines for presentation.

        Shared by the text, PDF and Office readers so that offset/limit,
        truncation and the continuation hints behave identically everywhere -
        previously each reader reimplemented this and they had drifted apart
        (only the text reader handled negative offsets correctly).

        :return: (result_dict, error_message) - exactly one is not None.
        """
        total_lines = len(all_lines)

        if total_lines == 0:
            return {
                "content": EMPTY_FILE_NOTICE,
                "total_lines": 0,
                "start_line": 0,
                "output_lines": 0,
                "is_empty": True,
            }, None

        start_line = 0
        if offset is not None:
            if offset < 0:
                # -20 means "the last 20 lines".
                start_line = max(0, total_lines + offset)
            else:
                start_line = max(0, offset - 1)
                if start_line >= total_lines:
                    return None, (
                        f"Error: Offset {offset} is beyond end of file "
                        f"({total_lines} lines total)"
                    )

        end_line = total_lines
        user_limited = False
        if limit is not None:
            end_line = min(start_line + limit, total_lines)
            user_limited = True

        selected = '\n'.join(all_lines[start_line:end_line])
        truncation = truncate_head(selected)
        start_display = start_line + 1

        if truncation.first_line_exceeds_limit:
            size = format_size(len(all_lines[start_line].encode('utf-8')))
            return {
                "content": (
                    f"[Line {start_display} is {size}, exceeds "
                    f"{format_size(DEFAULT_MAX_BYTES)} limit. Use bash tool to read: "
                    f"head -c {DEFAULT_MAX_BYTES} {display_path} | tail -n +{start_display}]"
                ),
                "total_lines": total_lines,
                "start_line": start_display,
                "output_lines": 0,
                "details": {"truncation": truncation.to_dict()},
            }, None

        output_text = number_lines(truncation.content, start_display)
        details = {}

        if truncation.truncated:
            end_display = start_display + truncation.output_lines - 1
            limit_note = (
                "" if truncation.truncated_by == "lines"
                else f" ({format_size(DEFAULT_MAX_BYTES)} limit)"
            )
            output_text += (
                f"\n\n[Showing lines {start_display}-{end_display} of {total_lines}"
                f"{limit_note}. Use offset={end_display + 1} to continue.]"
            )
            details["truncation"] = truncation.to_dict()
        elif user_limited and end_line < total_lines:
            output_text += (
                f"\n\n[{total_lines - end_line} more lines in file. "
                f"Use offset={end_line + 1} to continue.]"
            )

        result = {
            "content": output_text,
            "total_lines": total_lines,
            "start_line": start_display,
            "output_lines": truncation.output_lines,
        }
        if details:
            result["details"] = details
        return result, None

    def _is_credential_path(self, absolute_path: str) -> bool:
        """Shared guard; see agent.tools.utils.credentials.is_credential_path."""
        return is_credential_path(absolute_path)

    def _return_file_metadata(self, absolute_path: str, file_type: str, file_size: int) -> ToolResult:
        """
        Return file metadata for non-readable files (video, audio, binary, etc.)
        
        :param absolute_path: Absolute path to the file
        :param file_type: Type of file (video, audio, binary, etc.)
        :param file_size: File size in bytes
        :return: File metadata
        """
        file_name = Path(absolute_path).name
        file_ext = Path(absolute_path).suffix.lower()
        
        # Determine MIME type
        mime_types = {
            # Video
            '.mp4': 'video/mp4', '.avi': 'video/x-msvideo', '.mov': 'video/quicktime',
            '.mkv': 'video/x-matroska', '.webm': 'video/webm',
            # Audio
            '.mp3': 'audio/mpeg', '.wav': 'audio/wav', '.ogg': 'audio/ogg',
            '.m4a': 'audio/mp4', '.flac': 'audio/flac',
            # Binary
            '.zip': 'application/zip', '.tar': 'application/x-tar',
            '.gz': 'application/gzip', '.rar': 'application/x-rar-compressed',
        }
        mime_type = mime_types.get(file_ext, 'application/octet-stream')
        
        result = {
            "type": f"{file_type}_metadata",
            "file_type": file_type,
            "path": absolute_path,
            "file_name": file_name,
            "mime_type": mime_type,
            "size": file_size,
            "size_formatted": format_size(file_size),
            "message": f"{file_type.capitalize()} 文件: {file_name} ({format_size(file_size)})\n提示: 如果需要发送此文件，请使用 send 工具。"
        }
        
        return ToolResult.success(result)
    
    def _read_image(self, absolute_path: str, file_ext: str) -> ToolResult:
        """
        Read image file - always return metadata only (images should be sent, not read into context)
        
        :param absolute_path: Absolute path to the image file
        :param file_ext: File extension
        :return: Result containing image metadata for sending
        """
        try:
            # Get file size
            file_size = os.path.getsize(absolute_path)
            
            # Determine MIME type
            mime_type_map = {
                '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg',
                '.png': 'image/png',
                '.gif': 'image/gif',
                '.webp': 'image/webp'
            }
            mime_type = mime_type_map.get(file_ext, 'image/jpeg')
            
            # Return metadata for images (NOT file_to_send - use send tool to actually send)
            result = {
                "type": "image_metadata",
                "file_type": "image",
                "path": absolute_path,
                "mime_type": mime_type,
                "size": file_size,
                "size_formatted": format_size(file_size),
                "message": f"图片文件: {Path(absolute_path).name} ({format_size(file_size)})\n提示: 如果需要发送此图片，请使用 send 工具。"
            }
            
            return ToolResult.success(result)
            
        except Exception as e:
            return ToolResult.fail(f"Error reading image file: {str(e)}")
    
    def _read_text(self, absolute_path: str, display_path: str, offset: int = None, limit: int = None) -> ToolResult:
        """
        Read text file
        
        :param absolute_path: Absolute path to the file
        :param display_path: Path to display
        :param offset: Starting line number (1-indexed)
        :param limit: Maximum number of lines to read
        :return: File content or error message
        """
        try:
            # Check file size first
            file_size = os.path.getsize(absolute_path)
            MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
            
            if file_size > MAX_FILE_SIZE:
                # File too large, return metadata only
                return ToolResult.success({
                    "type": "file_to_send",
                    "file_type": "document",
                    "path": absolute_path,
                    "size": file_size,
                    "size_formatted": format_size(file_size),
                    "message": f"文件过大 ({format_size(file_size)} > 50MB)，无法读取内容。文件路径: {absolute_path}"
                })
            
            # Read file (utf-8-sig strips BOM automatically on Windows)
            # Note: Truncation is unified via truncate_head (DEFAULT_MAX_LINES / DEFAULT_MAX_BYTES)
            # so that offset/limit can paginate the entire file correctly.
            with open(absolute_path, 'r', encoding='utf-8-sig') as f:
                content = f.read()

            result, error = self._paginate(
                split_lines(content), offset, limit, display_path
            )
            if error:
                return ToolResult.fail(error)

            # Record the read so edit/write can warn if the file changes later.
            note_read(absolute_path)
            return ToolResult.success(result)
            
        except UnicodeDecodeError:
            return ToolResult.fail(f"Error: File is not a valid text file (encoding error): {display_path}")
        except Exception as e:
            return ToolResult.fail(f"Error reading file: {str(e)}")
    
    def _read_office(self, absolute_path: str, display_path: str, file_ext: str,
                     offset: int = None, limit: int = None) -> ToolResult:
        """Read Office documents (.docx, .xlsx, .pptx) using python-docx / openpyxl / python-pptx."""
        try:
            text = self._extract_office_text(absolute_path, file_ext)
        except ImportError as e:
            return ToolResult.fail(str(e))
        except Exception as e:
            return ToolResult.fail(f"Error reading Office document: {e}")

        if not text or not text.strip():
            return ToolResult.success({
                "content": f"[Office file {Path(absolute_path).name}: no text content could be extracted]",
            })

        result, error = self._paginate(
            split_lines(text), offset, limit, display_path
        )
        if error:
            return ToolResult.fail(error)

        note_read(absolute_path)
        return ToolResult.success(result)

    @staticmethod
    def _extract_office_text(absolute_path: str, file_ext: str) -> str:
        """Extract plain text from an Office document."""
        if file_ext in ('.docx', '.doc'):
            try:
                from docx import Document
            except ImportError:
                raise ImportError("Error: python-docx library not installed. Install with: pip install python-docx")
            doc = Document(absolute_path)
            paragraphs = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    paragraphs.append('\t'.join(cell.text for cell in row.cells))
            return '\n'.join(paragraphs)

        if file_ext in ('.xlsx', '.xls'):
            try:
                from openpyxl import load_workbook
            except ImportError:
                raise ImportError("Error: openpyxl library not installed. Install with: pip install openpyxl")
            wb = load_workbook(absolute_path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                parts.append(f"--- Sheet: {ws.title} ---")
                for row in ws.iter_rows(values_only=True):
                    parts.append('\t'.join(str(c) if c is not None else '' for c in row))
            wb.close()
            return '\n'.join(parts)

        if file_ext in ('.pptx', '.ppt'):
            try:
                from pptx import Presentation
            except ImportError:
                raise ImportError("Error: python-pptx library not installed. Install with: pip install python-pptx")
            prs = Presentation(absolute_path)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"--- Slide {i} ---")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                parts.append(text)
            return '\n'.join(parts)

        return ""

    def _read_pdf(self, absolute_path: str, display_path: str, offset: int = None,
                  limit: int = None, pages: str = None) -> ToolResult:
        """
        Read PDF file content
        
        :param absolute_path: Absolute path to the file
        :param display_path: Path to display
        :param offset: Starting line number (1-indexed)
        :param limit: Maximum number of lines to read
        :return: PDF text content or error message
        """
        try:
            # Try to import pypdf
            try:
                from pypdf import PdfReader
            except ImportError:
                return ToolResult.fail(
                    "Error: pypdf library not installed. Install with: pip install pypdf"
                )
            
            # Read PDF
            reader = PdfReader(absolute_path)
            total_pages = len(reader.pages)

            try:
                first_page, last_page = _parse_page_range(pages, total_pages)
            except ValueError as e:
                return ToolResult.fail(f"Error: {e}")

            # Only extract the requested window; a large PDF is expensive to
            # parse in full and the model rarely needs every page.
            text_parts = []
            for page_num in range(first_page, last_page + 1):
                page_text = reader.pages[page_num - 1].extract_text()
                if page_text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{page_text}")

            if not text_parts:
                scope = (
                    f"pages {first_page}-{last_page} of {total_pages}"
                    if (first_page, last_page) != (1, total_pages)
                    else f"{total_pages} pages"
                )
                return ToolResult.success({
                    "content": f"[PDF file with {scope}, but no text content could be extracted]",
                    "total_pages": total_pages,
                    "message": "PDF may contain only images or be encrypted"
                })

            result, error = self._paginate(
                split_lines("\n\n".join(text_parts)), offset, limit, display_path
            )
            if error:
                return ToolResult.fail(error)

            result["total_pages"] = total_pages
            result["pages_read"] = f"{first_page}-{last_page}"
            if last_page < total_pages:
                result["content"] += (
                    f"\n\n[Read pages {first_page}-{last_page} of {total_pages}. "
                    f"Use pages=\"{last_page + 1}-\" to continue.]"
                )

            note_read(absolute_path)
            return ToolResult.success(result)
            
        except Exception as e:
            return ToolResult.fail(f"Error reading PDF file: {str(e)}")
